from pathlib import Path
from pathlib import Path
# knowledge/file_processor.py (防崩溃版)
PROJECT_ROOT = Path(__file__).resolve().parent.parent
import os
from pathlib import Path

def extract_text_from_file(file_path: str) -> str:
    path = Path(file_path)
    if not path.exists():
        return ""
    suffix = path.suffix.lower()
    try:
        if suffix in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            return _safe_extract_image(file_path)
        elif suffix == '.pdf':
            return _safe_extract_pdf(file_path)
        elif suffix == '.docx':
            return _safe_extract_docx(file_path)
        elif suffix == '.pptx':
            return _safe_extract_pptx(file_path)
        elif suffix in ['.txt', '.md', '.py', '.json', '.yaml', '.csv']:
            return path.read_text(encoding='utf-8', errors='ignore')[:3000]
        else:
            # 其他文件尝试直接读文本，失败则返回空
            try:
                return path.read_text(encoding='utf-8', errors='ignore')[:1000]
            except:
                return ""
    except Exception as e:
        return f"[提取异常] {e}"

def _safe_extract_image(file_path):
    try:
        from PIL import Image
        img = Image.open(file_path)
        # 不调用 pytesseract 以免崩溃，只返回尺寸
        return f"[图片信息] 尺寸: {img.size}, 模式: {img.mode}"
    except Exception as e:
        return f"[图片读取失败] {e}"

def _safe_extract_pdf(file_path):
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text = "".join([page.extract_text() or "" for page in reader.pages])
        return text[:3000]
    except Exception as e:
        return f"[PDF提取失败] {e}"

def _safe_extract_docx(file_path):
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])[:3000]
    except Exception as e:
        return f"[DOCX提取失败] {e}"

def _safe_extract_pptx(file_path):
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)[:3000]
    except Exception as e:
        return f"[PPTX提取失败] {e}"
